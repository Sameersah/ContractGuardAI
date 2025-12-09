#!/usr/bin/env python3
"""
PPT Generator
Creates PowerPoint presentations from project information.
"""

import io
import logging
from typing import Optional, List, Tuple
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.warning("python-pptx not available. PPT files cannot be created.")

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


def create_project_presentation(output_path: str = "AWS_BOX_Project_Presentation.pptx") -> str:
    """
    Create a 5-6 slide PowerPoint presentation about the AWS-BOX Contract Protection System.
    
    Args:
        output_path: Path to save the PPTX file
        
    Returns:
        str: Path to the saved PPTX file
    """
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx is not installed. Please install it with: pip install python-pptx")
    
    try:
        # Create presentation
        prs = Presentation()
        
        # Set slide size to widescreen (16:9)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Slide 1: Title Slide
        _add_title_slide(prs)
        
        # Slide 2: Problem Statement
        _add_problem_slide(prs)
        
        # Slide 3: Solution Overview
        _add_solution_slide(prs)
        
        # Slide 4: Architecture
        _add_architecture_slide(prs)
        
        # Slide 5: Key Features
        _add_features_slide(prs)
        
        # Slide 6: Output & Benefits
        _add_benefits_slide(prs)
        
        # Save presentation
        prs.save(output_path)
        logger.info(f"PPT presentation saved to: {output_path}")
        return output_path
        
    except Exception as e:
        logger.error(f"Error creating PPT: {e}")
        raise


def _add_title_slide(prs):
    """Add title slide."""
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "AWS-BOX Contract Protection System"
    subtitle.text = "AI-Powered Contract Analysis & Protection\nAutomated Contract Review Using Box + AWS Bedrock"
    
    # Style the title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    subtitle.text_frame.paragraphs[0].font.size = Pt(20)


def _add_problem_slide(prs):
    """Add problem statement slide."""
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "The Problem"
    
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    problems = [
        "Contracts are complex and time-consuming to review",
        "Easy to miss unfavorable terms hidden in fine print",
        "No easy way to compare original vs. preferred terms",
        "Manual contract review requires legal expertise",
        "Multiple stakeholders need clarity on negotiation points"
    ]
    
    for i, problem in enumerate(problems):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = problem
        p.level = 0
        p.font.size = Pt(18)
        p.font.name = "Calibri"
        p.space_after = Pt(12)


def _add_solution_slide(prs):
    """Add solution overview slide."""
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "The Solution"
    
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    solutions = [
        "Automated contract monitoring via Box cloud storage",
        "AI-powered analysis using AWS Bedrock (Claude 3 Sonnet)",
        "Generates protected versions aligned with your interests",
        "Creates visual comparisons and negotiation guides",
        "Runs 24/7 on AWS EC2 with minimal infrastructure cost"
    ]
    
    for i, solution in enumerate(solutions):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = solution
        p.level = 0
        p.font.size = Pt(18)
        p.font.name = "Calibri"
        p.space_after = Pt(12)


def _add_architecture_slide(prs):
    """Add architecture slide."""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Title
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "System Architecture"
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    
    # Architecture components
    components = [
        ("Box API", Inches(1), Inches(2), Inches(2), Inches(1.5), "Monitors contracts folder"),
        ("AWS Bedrock", Inches(4), Inches(2), Inches(2), Inches(1.5), "AI Analysis Engine\n(Claude 3 Sonnet)"),
        ("EC2 Instance", Inches(7), Inches(2), Inches(2), Inches(1.5), "Contract Processor\n24/7 Monitoring"),
        ("Output Files", Inches(2.5), Inches(4.5), Inches(5), Inches(1.5), "3 Generated Documents\n• Mirror Contract\n• Redline Comparison\n• Negotiation Guide")
    ]
    
    for label, left_pos, top_pos, width_pos, height_pos, desc in components:
        # Add rectangle
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_pos, top_pos, width_pos, height_pos)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(31, 78, 121)  # Blue
        shape.line.color.rgb = RGBColor(0, 0, 0)
        
        # Add text
        text_frame = shape.text_frame
        text_frame.text = f"{label}\n{desc}"
        text_frame.word_wrap = True
        text_frame.paragraphs[0].font.size = Pt(14)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Center text vertically
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def _add_features_slide(prs):
    """Add key features slide."""
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Key Features"
    
    text_frame = body_shape.text_frame
    text_frame.clear()
    
    features = [
        "Real-time contract detection and processing",
        "Customizable interest profiles (MY_INTERESTS.txt)",
        "Per-contract instruction overrides",
        "Three comprehensive output documents per contract",
        "Automated email notifications via AWS SNS",
        "Secure credential management via AWS Secrets Manager",
        "Cost-effective: ~$8-11/month on AWS"
    ]
    
    for i, feature in enumerate(features):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = feature
        p.level = 0
        p.font.size = Pt(16)
        p.font.name = "Calibri"
        p.space_after = Pt(10)


def _add_benefits_slide(prs):
    """Add output and benefits slide."""
    two_content_layout = prs.slide_layouts[3]
    slide = prs.slides.add_slide(two_content_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    left_shape = shapes.placeholders[1]
    right_shape = shapes.placeholders[2]
    
    title_shape.text = "Output & Benefits"
    
    # Left side - Outputs
    left_frame = left_shape.text_frame
    left_frame.clear()
    left_frame.text = "Generated Documents"
    left_frame.paragraphs[0].font.size = Pt(20)
    left_frame.paragraphs[0].font.bold = True
    
    outputs = [
        "1. Mirror Contract\n   Protected version with your interests",
        "2. Redline Comparison\n   Visual diff showing all changes",
        "3. Negotiation Guide\n   How to present your terms"
    ]
    
    for i, output in enumerate(outputs):
        p = left_frame.add_paragraph()
        p.text = output
        p.level = 0
        p.font.size = Pt(14)
        p.space_after = Pt(12)
    
    # Right side - Benefits
    right_frame = right_shape.text_frame
    right_frame.clear()
    right_frame.text = "Key Benefits"
    right_frame.paragraphs[0].font.size = Pt(20)
    right_frame.paragraphs[0].font.bold = True
    
    benefits = [
        "⚡ Saves hours of manual review",
        "🎯 Ensures your interests are protected",
        "📊 Clear visualization of changes",
        "🤝 Better negotiation preparation",
        "🔒 Secure and automated",
        "💰 Cost-effective solution"
    ]
    
    for i, benefit in enumerate(benefits):
        p = right_frame.add_paragraph()
        p.text = benefit
        p.level = 0
        p.font.size = Pt(14)
        p.space_after = Pt(10)


if __name__ == "__main__":
    import sys
    
    output_file = sys.argv[1] if len(sys.argv) > 1 else "AWS_BOX_Project_Presentation.pptx"
    
    try:
        output_path = create_project_presentation(output_file)
        print(f"✓ PowerPoint presentation created: {output_path}")
    except Exception as e:
        print(f"✗ Error: {e}")
        sys.exit(1)

